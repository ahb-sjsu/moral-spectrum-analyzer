"""Build the GTC Information Capsule as a self-running, narrated 16:9 PPTX.

Renders the 8-shot storyboard in docs/CAPSULE.md: dark deck, the recurring
moral-spectrum bar motif, on-slide text, the voiceover in speaker notes, a fade
transition + per-shot auto-advance timing (so it self-runs to ~90 s). PowerPoint
entrance animations (bar-grow, badge-pop, typewriter) are added in-app; each
slide's speaker notes carry its animation cue.

    python scripts/build_capsule_pptx.py            # -> out/capsule.pptx
Export the video in PowerPoint: File > Export > Create a Video (use recorded
timings/narration) -> capsule.mp4.
"""
from __future__ import annotations

import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---- palette ----------------------------------------------------------------
BG = RGBColor(0x0B, 0x10, 0x20)
WHITE = RGBColor(0xF2, 0xF5, 0xFA)
MUTE = RGBColor(0x8A, 0x93, 0xA8)
GREEN = RGBColor(0x3D, 0xDC, 0x97)   # validated / accent
GOLD = RGBColor(0xF5, 0xC1, 0x47)    # brand gold (the Eris apple) — for the slogan
AMBER = RGBColor(0xF2, 0xB1, 0x3D)   # escalate
RED = RGBColor(0xE5, 0x5A, 0x5A)     # remove
RAMP = [RGBColor(*c) for c in [
    (0x4C, 0x8B, 0xF5), (0x53, 0xA0, 0xEE), (0x5B, 0xB8, 0xD8), (0x62, 0xCB, 0xB0),
    (0x6E, 0xD1, 0x8C), (0x9C, 0xD1, 0x6E), (0xCF, 0xC7, 0x55), (0xEC, 0xA8, 0x4C),
    (0xE8, 0x82, 0x50)]]
IDENT = RGBColor(0xB8, 0x6B, 0xF5)   # identity-attack (violet, the discovered axis)

# short axis labels (canonical DEME9 + discovered 10th)
AX9 = ["harm", "rights", "fairness", "autonomy", "privacy",
       "environ", "care", "legitimacy", "epistemic"]
AX10 = AX9 + ["identity‑attack"]

# per-shot advance seconds (sums to 90)
TIMING = [11, 12, 12, 14, 12, 13, 10, 6]
INTRO_SEC = 5   # silent branded title card before the narration begins

EMU_IN = 914400


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BG


def text(slide, s, l, t, w, h, size=28, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Segoe UI"
    return tb


def rect(slide, l, t, w, h, color, line=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = color
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def chip(slide, s, l, t, w, color, fg=WHITE, size=20):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(0.55), )
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = fg; r.font.name = "Segoe UI"
    return sp


def bars(slide, heights, labels, highlight=None, badges=False,
         base_y=6.1, max_h=3.2, left=1.15, right=12.2):
    n = len(heights)
    span = right - left
    slot = span / n
    bw = slot * 0.62
    for i, (hgt, lab) in enumerate(zip(heights, labels)):
        x = left + i * slot + (slot - bw) / 2
        h = max(0.06, hgt * max_h)
        col = IDENT if (labels[i].startswith("identity")) else RAMP[i % len(RAMP)]
        if highlight is not None and i != highlight and highlight >= 0:
            pass
        rect(slide, x, base_y - h, bw, h, col)
        if badges:
            text(slide, "✓", x, base_y - h - 0.42, bw, 0.4, size=16,
                 color=GREEN, align=PP_ALIGN.CENTER)
        text(slide, lab, x - slot * 0.15, base_y + 0.06, slot * 0.9, 0.5,
             size=10.5, color=MUTE, align=PP_ALIGN.CENTER)


def notes(slide, vo, cue):
    tf = slide.notes_slide.notes_text_frame
    tf.text = f"VOICEOVER:\n{vo}\n\nANIMATION CUE:\n{cue}"


def advance(slide, seconds):
    """Fade transition + auto-advance after `seconds` (click also advances)."""
    sld = slide._element
    for old in sld.findall(qn("p:transition")):
        sld.remove(old)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", "med")
    tr.set("advClick", "1")
    tr.set("advTm", str(int(seconds * 1000)))
    etree.SubElement(tr, qn("p:fade"))


# ---- golden apple (ErisML / Eris mark) --------------------------------------
# Mirrors docs/brand/erisml_apple.svg exactly (path coords + the load-bearing
# brand colours). Rendered with Pillow because no SVG rasteriser (cairo) is
# available here; keep in sync with the SVG if that master changes. Do NOT
# recolour — the gold is the Eris reference (brand kit rule).
_APPLE_PNG = None


def _quad(p0, p1, p2, n=48):
    return [((1 - t) ** 2 * p0[0] + 2 * (1 - t) * t * p1[0] + t * t * p2[0],
             (1 - t) ** 2 * p0[1] + 2 * (1 - t) * t * p1[1] + t * t * p2[1])
            for t in (i / n for i in range(n + 1))]


def render_apple(path, px=560, ss=3):
    from PIL import Image, ImageDraw
    W = px * ss
    S = W / 200.0

    def sc(p):
        return (p[0] * S, p[1] * S)

    GOLD = (0xF5, 0xC1, 0x47, 255)
    HI = (0xFF, 0xE8, 0xA8, 128)      # highlight, opacity 0.5
    STEM = (0x6B, 0x3A, 0x00, 255)
    LEAF = (0x4A, 0xDE, 0x80, 255)
    VEIN = (0x16, 0x65, 0x34, 255)
    img = Image.new("RGBA", (W, W), (0, 0, 0, 0))
    d = ImageDraw.Draw(img, "RGBA")
    d.ellipse([sc((42, 57)), sc((158, 173))], fill=GOLD)          # body: cx100 cy115 r58
    hi = Image.new("RGBA", (W, W), (0, 0, 0, 0))
    ImageDraw.Draw(hi).ellipse([sc((70, 75)), sc((98, 115))], fill=HI)  # cx84 cy95 rx14 ry20
    img = Image.alpha_composite(img, hi)
    d = ImageDraw.Draw(img, "RGBA")
    d.line([sc(p) for p in _quad((96, 60), (100, 46), (112, 50))],       # stem
           fill=STEM, width=max(2, int(3.5 * S)), joint="curve")
    leaf = _quad((114, 50), (135, 42), (140, 62)) + _quad((140, 62), (124, 67), (114, 50))
    d.polygon([sc(p) for p in leaf], fill=LEAF)                          # leaf
    d.line([sc(p) for p in _quad((114, 50), (130, 52), (140, 62))],      # vein
           fill=VEIN, width=max(1, int(1.2 * S)), joint="curve")
    img.resize((px, px), Image.LANCZOS).save(path)


def apple_png():
    global _APPLE_PNG
    if _APPLE_PNG is None:
        os.makedirs("out", exist_ok=True)
        _APPLE_PNG = os.path.join("out", "_apple.png")
        render_apple(_APPLE_PNG)
    return _APPLE_PNG


def add_apple(slide, l, t, size):
    slide.shapes.add_picture(apple_png(), Inches(l), Inches(t),
                             width=Inches(size), height=Inches(size))


# ---- analyzer sub-mark (apple -> refracted moral spectrum) -------------------
# Mirrors docs/brand/analyzer_mark.svg exactly (viewBox 400x200): the apple
# embedded unchanged, refracting into the nine-axis spectrum. Transparent PNG so
# it sits on the deck background. Keep in sync with the SVG master.
_MARK_PNG = None
_MARK_H = [0.45, 0.60, 0.78, 0.68, 0.90, 0.75, 0.85, 0.60, 0.72]
_MARK_RAMP = [(0xF5, 0xC1, 0x47), (0xEC, 0xA8, 0x4C), (0xCF, 0xC7, 0x55),
              (0x9C, 0xD1, 0x6E), (0x6E, 0xD1, 0x8C), (0x62, 0xCB, 0xB0),
              (0x5B, 0xB8, 0xD8), (0x53, 0xA0, 0xEE), (0x4C, 0x8B, 0xF5)]


def render_mark(path, w_px=1600, ss=2):
    from PIL import Image, ImageDraw
    W, Hp = w_px * ss, (w_px // 2) * ss
    S = W / 400.0
    base, maxh, left, right, bw = 155.0, 95.0, 176.0, 380.0, 15.0
    n = len(_MARK_H); slot = (right - left) / n
    edge = (131.0, 104.0)
    img = Image.new("RGBA", (W, Hp), (0, 0, 0, 0))
    d = ImageDraw.Draw(img, "RGBA")

    def M(p):                                   # mark-space -> px
        return (p[0] * S, p[1] * S)

    def A(p):                                   # apple-space -> mark-space -> px
        return ((p[0] * 0.75 + 15) * S, (p[1] * 0.75 + 18.75) * S)
    # refraction rays: apple edge -> each bar top
    for i, h in enumerate(_MARK_H):
        cx = left + i * slot + slot / 2
        top = base - h * maxh
        d.line([M(edge), M((cx, top))], fill=_MARK_RAMP[i] + (150,),
               width=max(1, int(1.3 * S)))
    # spectrum bars
    for i, h in enumerate(_MARK_H):
        x = left + i * slot + (slot - bw) / 2
        top = base - h * maxh
        d.rounded_rectangle([M((x, top)), M((x + bw, base))], radius=int(2 * S),
                            fill=_MARK_RAMP[i] + (255,))
    # apple (embedded unchanged), scaled/placed via A()
    d.ellipse([A((42, 57)), A((158, 173))], fill=(0xF5, 0xC1, 0x47, 255))
    hi = Image.new("RGBA", (W, Hp), (0, 0, 0, 0))
    ImageDraw.Draw(hi).ellipse([A((70, 75)), A((98, 115))], fill=(0xFF, 0xE8, 0xA8, 128))
    img = Image.alpha_composite(img, hi); d = ImageDraw.Draw(img, "RGBA")
    d.line([A(p) for p in _quad((96, 60), (100, 46), (112, 50))],
           fill=(0x6B, 0x3A, 0x00, 255), width=max(2, int(3.5 * 0.75 * S)), joint="curve")
    leaf = _quad((114, 50), (135, 42), (140, 62)) + _quad((140, 62), (124, 67), (114, 50))
    d.polygon([A(p) for p in leaf], fill=(0x4A, 0xDE, 0x80, 255))
    d.line([A(p) for p in _quad((114, 50), (130, 52), (140, 62))],
           fill=(0x16, 0x65, 0x34, 255), width=max(1, int(1.2 * 0.75 * S)), joint="curve")
    img.resize((w_px, w_px // 2), Image.LANCZOS).save(path)


def mark_png():
    global _MARK_PNG
    if _MARK_PNG is None:
        os.makedirs("out", exist_ok=True)
        _MARK_PNG = os.path.join("out", "_mark.png")
        render_mark(_MARK_PNG)
    return _MARK_PNG


def add_mark(slide, l, t, w):
    slide.shapes.add_picture(mark_png(), Inches(l), Inches(t),
                             width=Inches(w), height=Inches(w / 2))


def ray(slide, x1, y1, x2, y2, color, width=0.75):
    """A thin refraction ray (the apple's light fanning into the spectrum)."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def build(path="out/capsule.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    VO = [
        "Every day, automated systems decide what billions of people can and can't say. "
        "Most reduce that judgment — a moral one — to a single number. A toxicity score.",
        "But one number can't tell you which value it acted on. It shifts when you reword the "
        "same sentence. And worst of all — it can't tell you what it fails to see.",
        "So we built the Moral Spectrum Analyzer. It reads content not as one score, but as a "
        "spectrum — the energy across nine distinct moral dimensions.",
        "Feed it a line of content. Each axis lights up — false information, real-world harm, "
        "eroded public trust — and every one carries a badge showing it was independently "
        "validated. No badge, no number.",
        "Out comes a decision — allow, remove, or escalate to a human — with its moral "
        "residue: every value it weighed. Here it escalates: the disinformation signal is loud, "
        "but the call goes to a person. Nothing hidden.",
        "Nine of ten axes pass a pre-registered validation gate. Reword the sentence, or translate "
        "it across five languages — the verdict holds, where a single score would drift. And "
        "every decision emits a proof you can re-verify.",
        "It even discovered a moral dimension the standard taxonomy was missing — identity "
        "attack — and validated it. And when it can't read something? It says so.",
        "The Moral Spectrum Analyzer. Trust — made measurable.",
    ]

    # --- Intro: branded title card (silent cold open, holds then narration begins) ---
    s = _blank(prs); bg(s)
    add_mark(s, (13.333 - 7.4) / 2, 0.85, 7.4)          # analyzer mark: apple -> spectrum
    text(s, "Moral Spectrum Analyzer", 0, 4.75, 13.333, 1.0, size=46, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "a judgment isn't one score — it has a spectrum", 0, 5.72, 13.333, 0.6,
         size=22, color=GREEN, align=PP_ALIGN.CENTER, italic=True)
    text(s, "An Information Capsule  ·  a Philosophy Engineering discipline", 0, 6.42,
         13.333, 0.5, size=17, color=GOLD, align=PP_ALIGN.CENTER)
    text(s, "Global Trust Challenge 2026  ·  Geometric Ethics AI Lab", 0, 6.86, 13.333,
         0.5, size=14, color=MUTE, align=PP_ALIGN.CENTER)
    notes(s, "(silent title card — no narration; holds ~5 s, then Shot 1 begins)",
          "Mark fades up; spectrum bars settle under the apple; title + tagline rise.")
    advance(s, INTRO_SEC)

    # --- Shot 1: the lone number ---
    s = _blank(prs); bg(s)
    text(s, "0.82", 0, 2.1, 13.333, 2.0, size=140, color=MUTE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, "one number.", 0, 4.5, 13.333, 1.0, size=34, color=WHITE,
         align=PP_ALIGN.CENTER, italic=True)
    text(s, "AI decides what billions can and can't say — as a single toxicity score.",
         0, 5.6, 13.333, 0.8, size=18, color=MUTE, align=PP_ALIGN.CENTER)
    notes(s, VO[0], "Number pulses once; blurred feed of posts drifts up behind it.")
    advance(s, TIMING[0])

    # --- Shot 2: three failures ---
    s = _blank(prs); bg(s)
    text(s, "one number, three blind spots", 0, 0.7, 13.333, 0.9, size=32,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    cols = [("⊘", "OPAQUE", "which value did it act on?"),
            ("⇄", "UNSTABLE", "reword it — the score moves"),
            ("◑", "BLIND", "it can't say what it can't see")]
    for i, (ic, hd, sub) in enumerate(cols):
        x = 1.0 + i * 3.95
        text(s, ic, x, 2.2, 3.4, 1.2, size=64, color=AMBER, align=PP_ALIGN.CENTER)
        text(s, hd, x, 3.6, 3.4, 0.7, size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(s, sub, x, 4.3, 3.4, 1.0, size=17, color=MUTE, align=PP_ALIGN.CENTER)
    notes(s, VO[1], "Each fault fades in on its VO beat (opaque, unstable, blind).")
    advance(s, TIMING[1])

    # --- Shot 3: the reveal — the apple refracts into its moral spectrum ---
    s = _blank(prs); bg(s)
    text(s, "Moral Spectrum Analyzer", 0, 0.55, 13.333, 0.9, size=38, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "a judgment isn't one score — it has a spectrum",
         0, 1.45, 13.333, 0.6, size=20, color=GREEN, align=PP_ALIGN.CENTER, italic=True)
    ax, ay, asz = 1.35, 3.7, 1.75          # apple (the judgment) on the left
    add_apple(s, ax, ay - asz / 2, asz)
    hs = [0.5, 0.35, 0.55, 0.3, 0.4, 0.25, 0.6, 0.35, 0.45]
    left, right, base_y, max_h = 4.7, 12.3, 6.1, 2.7
    n = len(hs); slot = (right - left) / n
    apex = (ax + asz - 0.2, ay)            # refraction origin: the apple's edge
    for i, h in enumerate(hs):             # rays fan from the apple to each bar top
        ray(s, apex[0], apex[1], left + i * slot + slot / 2, base_y - h * max_h,
            RAMP[i % len(RAMP)], 0.75)
    bars(s, hs, AX9, left=left, right=right, base_y=base_y, max_h=max_h)
    notes(s, VO[2], "The apple's golden light refracts, prism-like, into the nine spectrum "
          "bars: rays fan from the apple to each bar, then bars grow L->R.")
    advance(s, TIMING[2])

    # --- Shot 4: walkthrough input -> process ---
    s = _blank(prs); bg(s)
    text(s, "input", 1.15, 0.55, 3, 0.5, size=18, color=MUTE)
    rect(s, 1.15, 1.05, 11.05, 0.9, RGBColor(0x16, 0x1E, 0x33))
    text(s, "“Doctors are hiding the cure — drink bleach to flush the virus.”",
         1.45, 1.15, 10.5, 0.7, size=24, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "each axis validated — no badge, no number", 1.15, 2.15, 11, 0.5,
         size=15, color=GREEN, italic=True)
    # disinformation: epistemic (false info) + harm + legitimacy dominate
    bars(s, [0.75, 0.15, 0.15, 0.35, 0.08, 0.05, 0.25, 0.65, 0.92], AX9,
         badges=True, base_y=6.2, max_h=2.9)
    notes(s, VO[3], "Typewriter the input; bars ease to height; green ✓ badges pop in staggered.")
    advance(s, TIMING[3])

    # --- Shot 5: output decision + residue ---
    s = _blank(prs); bg(s)
    text(s, "output", 1.15, 0.55, 3, 0.5, size=18, color=MUTE)
    chip(s, "ESCALATE → human", 1.15, 1.15, 4.2, AMBER, fg=RGBColor(0x11, 0x11, 0x11))
    text(s, "moral residue — every value weighed:", 6.0, 1.15, 6.2, 0.5,
         size=16, color=MUTE)
    text(s, "epistemic · harm · legitimacy", 6.0, 1.65, 6.2, 0.6,
         size=18, color=WHITE)
    text(s, "allow  ·  remove  ·  escalate — the call goes to a person, nothing hidden",
         1.15, 2.15, 11, 0.5, size=15, color=MUTE, italic=True)
    bars(s, [0.75, 0.15, 0.15, 0.35, 0.08, 0.05, 0.25, 0.65, 0.92], AX9,
         base_y=6.2, max_h=2.9)
    notes(s, VO[4], "Decision chip snaps in; residue panel slides in from the right.")
    advance(s, TIMING[4])

    # --- Shot 6: grounded / invariant / auditable triptych ---
    s = _blank(prs); bg(s)
    text(s, "grounded · invariant · auditable", 0, 0.6, 13.333, 0.8, size=30,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    panels = [
        ("9 / 10", "axes pass a pre-registered\nvalidation gate", GREEN),
        ("es ar zh hi sw", "reword or translate —\nthe verdict holds", RGBColor(0x5B, 0xB8, 0xD8)),
        ("✓ re-verified", "every decision emits a\nhash-chained proof", GREEN),
    ]
    for i, (big, sub, col) in enumerate(panels):
        x = 0.7 + i * 4.1
        rect(s, x, 2.0, 3.7, 3.0, RGBColor(0x14, 0x1C, 0x30))
        text(s, big, x, 2.4, 3.7, 1.1, size=40, color=col, bold=True, align=PP_ALIGN.CENTER)
        text(s, sub, x, 3.7, 3.7, 1.1, size=18, color=WHITE, align=PP_ALIGN.CENTER)
    notes(s, VO[5], "Three panels wipe in L->R on the VO beats (gate stamp, languages collapse to one point, audit check).")
    advance(s, TIMING[5])

    # --- Shot 7: discovery + escalate ---
    s = _blank(prs); bg(s)
    text(s, "it finds what the taxonomy missed", 0, 0.7, 13.333, 0.8, size=32,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    bars(s, [0.5, 0.35, 0.55, 0.3, 0.4, 0.25, 0.6, 0.35, 0.45, 0.85], AX10,
         base_y=5.9, max_h=2.7)
    chip(s, "identity attack   ✓ validated · AUROC 0.80", 3.55, 6.35, 6.25, IDENT)
    notes(s, VO[6], "A NEW violet 10th bar rises out of the old blind-spot wedge + badge; hold. (Escalation was already shown in Shot 5.)")
    advance(s, TIMING[6])

    # --- Shot 8: golden-apple lockup + slogan ---
    s = _blank(prs); bg(s)
    add_apple(s, (13.333 - 1.5) / 2, 0.8, 1.5)   # hero apple, centered
    text(s, "Moral Spectrum Analyzer", 0, 2.7, 13.333, 1.0, size=44, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "Trust — made measurable.", 0, 3.8, 13.333, 0.9, size=30, color=GOLD,
         align=PP_ALIGN.CENTER, italic=True)
    bars(s, [0.12] * 9, [""] * 9, base_y=5.25, max_h=0.5)   # quiet spectrum underline
    text(s, "a Philosophy Engineering discipline", 0, 5.95, 13.333, 0.5, size=17,
         color=GOLD, align=PP_ALIGN.CENTER)
    text(s, "Global Trust Challenge 2026  ·  Geometric Ethics AI Lab",
         0, 6.4, 13.333, 0.6, size=15, color=MUTE, align=PP_ALIGN.CENTER)
    notes(s, VO[7], "Golden apple settles in; wordmark + slogan fade up; bars sweep once beneath.")
    advance(s, TIMING[7])

    # subtle golden-apple brand mark, top-right, on content slides that don't already
    # feature the apple (skip the intro title card, Shot 3's refraction reveal, and
    # Shot 8's hero lockup — indices shift +1 because of the prepended intro)
    for i, sl in enumerate(prs.slides):
        if i not in (0, 3, 8):
            add_apple(sl, 12.55, 0.32, 0.5)

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    total = INTRO_SEC + sum(TIMING)
    print(f"wrote {path}  ({len(list(prs.slides))} slides, ~{total}s auto-run "
          f"[{INTRO_SEC}s intro + {sum(TIMING)}s narrated])")


if __name__ == "__main__":
    build()
